"""Ingestion: page specs, text/Office loaders, mode routing."""
import pytest

from diannot.ingest import (
    extract_text_from_docx,
    extract_text_from_pptx,
    load_raw_text,
    parse_pages,
)
from diannot.pipeline import decide_mode


def test_parse_pages():
    assert parse_pages("1-3,5", 10) == [0, 1, 2, 4]
    assert parse_pages("5-2", 10) == [1, 2, 3, 4]  # reversed range tolerated
    assert parse_pages("2", 5) == [1]
    with pytest.raises(ValueError):
        parse_pages("99", 5)  # out of range -> nothing selected


def test_load_txt_and_unsupported(tmp_path):
    p = tmp_path / "n.txt"
    p.write_text("hello", encoding="utf-8")
    assert load_raw_text(p) == "hello"
    bad = tmp_path / "n.xyz"
    bad.write_text("x", encoding="utf-8")
    with pytest.raises(ValueError):
        load_raw_text(bad)


def test_docx(tmp_path):
    import docx

    d = docx.Document()
    d.add_paragraph("The heart pumps blood.")
    f = tmp_path / "a.docx"
    d.save(f)
    assert "heart" in extract_text_from_docx(f).lower()


def test_pptx(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Vessels"
    f = tmp_path / "a.pptx"
    prs.save(f)
    assert "Vessels" in extract_text_from_pptx(f)


def test_decide_mode_image_and_text():
    assert decide_mode(".png", None, False, "x.png", None) == "vision"
    assert decide_mode(".png", None, True, "x.png", None) == "tesseract"
    assert decide_mode(".png", False, False, "x.png", None) == "tesseract"
    assert decide_mode(".txt", None, False, "x.txt", None) == "text"
    assert decide_mode(".docx", None, False, "x.docx", None) == "text"


def test_decide_mode_pdf(tmp_path):
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "This text PDF has plenty of real selectable text content.")
    f = tmp_path / "t.pdf"
    doc.save(f)
    doc.close()
    assert decide_mode(".pdf", None, False, f, None) == "text"  # has text -> not scanned
    assert decide_mode(".pdf", True, False, f, None) == "vision"  # forced
