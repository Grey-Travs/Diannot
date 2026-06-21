"""Load raw text from supported inputs.

Plain text (``.txt``/``.md``) and *text* PDFs are read here via PyMuPDF. Images and scanned PDFs
are routed to vision (or Tesseract) by :mod:`diannot.pipeline`; Word/PowerPoint are handled there too.
"""
from __future__ import annotations

from pathlib import Path

TEXT_SUFFIXES = {".txt", ".md", ".text"}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff"}


def parse_pages(spec: str, total: int) -> list[int]:
    """Parse a 1-based page spec like ``"1-3,5"`` into sorted 0-based indices."""
    indices: set[int] = set()
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            lo_s, hi_s = part.split("-", 1)
            lo, hi = int(lo_s), int(hi_s)
            if lo > hi:  # tolerate reversed ranges like "5-2"
                lo, hi = hi, lo
            for n in range(lo, hi + 1):
                if 1 <= n <= total:
                    indices.add(n - 1)
        else:
            n = int(part)
            if 1 <= n <= total:
                indices.add(n - 1)
    if not indices:
        raise ValueError(f"Page spec '{spec}' selected no pages (document has {total}).")
    return sorted(indices)


def extract_text_from_pdf(path: Path | str, pages: str | None = None) -> str:
    """Extract plain text from a PDF (optionally a subset of pages)."""
    import fitz  # PyMuPDF

    path = Path(path)
    doc = fitz.open(path)
    try:
        idxs = parse_pages(pages, doc.page_count) if pages else range(doc.page_count)
        parts = [doc[i].get_text() for i in idxs]
    finally:
        doc.close()
    return "\n\n".join(parts).strip()


def extract_text_from_docx(path: Path | str) -> str:
    """Extract text (paragraphs + tables) from a Word ``.docx`` file."""
    import docx

    document = docx.Document(str(path))
    parts: list[str] = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def extract_text_from_pptx(path: Path | str) -> str:
    """Extract text from a PowerPoint ``.pptx`` (one block per slide)."""
    from pptx import Presentation

    prs = Presentation(str(path))
    slides: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        lines.append(" | ".join(cells))
        if lines:
            slides.append(f"# Slide {index}\n" + "\n".join(lines))
    return "\n\n".join(slides).strip()


def load_raw_text(path: Path | str, pages: str | None = None) -> str:
    """Load raw text from a ``.txt``/``.md``, ``.pdf``, ``.docx`` or ``.pptx`` file."""
    path = Path(path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_text_from_pdf(path, pages)
    if suffix == ".docx":
        return extract_text_from_docx(path)
    if suffix == ".pptx":
        return extract_text_from_pptx(path)
    if suffix in TEXT_SUFFIXES:
        return path.read_text(encoding="utf-8").strip()
    raise ValueError(
        f"Unsupported text input '{path.name}' (expected .txt, .md, .pdf, .docx or .pptx)."
    )


def is_scanned_pdf(path: Path | str, pages: str | None = None, min_chars_per_page: int = 20) -> bool:
    """Heuristic: a PDF with almost no extractable text is image-only (scanned)."""
    import fitz

    doc = fitz.open(path)
    try:
        # Sample the first few pages when no range is given — a scanned PDF is scanned throughout,
        # so this detects it without reading every page (which would stall the import UI on a big PDF).
        idxs = list(parse_pages(pages, doc.page_count) if pages else range(min(doc.page_count, 5)))
        chars = sum(len(doc[i].get_text().strip()) for i in idxs)
    finally:
        doc.close()
    return chars < min_chars_per_page * max(len(idxs), 1)


def load_image_sources(path: Path | str, pages: str | None = None, dpi: int = 200) -> list[bytes]:
    """Return page image(s) as PNG bytes — for an image file or a (rasterized) PDF."""
    import fitz

    path = Path(path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        doc = fitz.open(path)
        try:
            idxs = parse_pages(pages, doc.page_count) if pages else range(doc.page_count)
            return [doc[i].get_pixmap(dpi=dpi).tobytes("png") for i in idxs]
        finally:
            doc.close()
    if suffix in IMAGE_SUFFIXES:
        pix = fitz.Pixmap(str(path))
        if pix.colorspace is not None and pix.colorspace.n >= 4:  # CMYK -> RGB
            pix = fitz.Pixmap(fitz.csRGB, pix)
        return [pix.tobytes("png")]
    raise ValueError(f"'{path.name}' is not an image or PDF.")


def page_numbers_for(path: Path | str, pages: str | None = None) -> list[int]:
    """1-based page numbers that :func:`load_image_sources` would produce, in order."""
    path = Path(path)
    if path.suffix.lower() == ".pdf":
        import fitz

        doc = fitz.open(path)
        try:
            idxs = parse_pages(pages, doc.page_count) if pages else range(doc.page_count)
            return [i + 1 for i in idxs]
        finally:
            doc.close()
    return [1]  # a single image file is "page 1"


def ocr_image_sources(images: list[bytes], lang: str = "eng") -> str:
    """Offline OCR of PNG image bytes via Tesseract (the `ocr` extra)."""
    try:
        import io

        import pytesseract
        from PIL import Image
    except ImportError as exc:  # pragma: no cover - optional dependency
        raise RuntimeError(
            "Offline OCR needs the 'ocr' extra (uv sync --extra ocr) and the Tesseract "
            "binary installed and on PATH."
        ) from exc
    parts = [pytesseract.image_to_string(Image.open(io.BytesIO(b)), lang=lang) for b in images]
    return "\n\n".join(parts).strip()
